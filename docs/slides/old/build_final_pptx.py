from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
SLIDES_DIR = ROOT / "docs" / "slides"
FIG_DIR = ROOT / "docs" / "figs"
TEMPLATE = Path(
    "/Users/kfunaki/Documents/1.C_NYU/VIP_Spring_2026/"
    "VIP_Spring_2026_April_kf2623.pptx"
)
OUT = SLIDES_DIR / "vip_final_5_slides_template.pptx"
NOTES_OUT = SLIDES_DIR / "vip_final_5_slides_template_speaker_notes.md"

BLUE = RGBColor(54, 96, 145)
DARK = RGBColor(45, 45, 45)
GRAY = RGBColor(105, 105, 105)
LIGHT_BLUE = RGBColor(225, 235, 247)
LIGHT_GRAY = RGBColor(242, 242, 242)
GREEN = RGBColor(91, 155, 104)
RED = RGBColor(178, 76, 76)
WHITE = RGBColor(255, 255, 255)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[index]


def clear_template_slides(prs: Presentation) -> None:
    for idx in range(len(prs.slides) - 1, -1, -1):
        delete_slide(prs, idx)


def add_footer(slide, page: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(5.32), Inches(9.3), Inches(0.18))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"VIP Forecasting Project | Kenji Funaki | May 2026 | {page}/5"
    p.font.name = "Arial"
    p.font.size = Pt(7)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title: str, page: int) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.10))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.20), Inches(9.2), Inches(0.42))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE
    add_footer(slide, page)


def add_label(slide, text: str, x, y, w, h, fill=BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_textbox(slide, text: str, x, y, w, h, size=11, bold=False, color=DARK) -> None:
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color


def add_bullets(slide, bullets, x, y, w, h, size=11, color=DARK) -> None:
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)


def style_cell(cell, text, size=9, bold=False, color=DARK, fill=None, align=PP_ALIGN.CENTER):
    cell.text = text
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_table(slide, rows, x, y, w, h, font_size=8, first_col_left=True):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            fill = BLUE if r == 0 else (LIGHT_BLUE if r % 2 == 1 else WHITE)
            color = WHITE if r == 0 else DARK
            align = PP_ALIGN.LEFT if first_col_left and c == 0 and r > 0 else PP_ALIGN.CENTER
            style_cell(table.cell(r, c), value, size=font_size, bold=(r == 0), color=color, fill=fill, align=align)
    return table_shape


def add_callout(slide, text, x, y, w, h, fill=LIGHT_BLUE, line=BLUE, size=10) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.06)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.color.rgb = DARK


def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "1. Alternative Data Update: Narrower, Cleaner, More Honest", 1)
    add_label(slide, "April issue", Inches(0.45), Inches(0.85), Inches(2.6), Inches(0.35), RED)
    add_label(slide, "May update", Inches(3.70), Inches(0.85), Inches(2.6), Inches(0.35), BLUE)
    add_label(slide, "Interpretation", Inches(6.95), Inches(0.85), Inches(2.6), Inches(0.35), GREEN)
    add_bullets(
        slide,
        [
            "Broad alternative-data list was too noisy for daily direction.",
            "Validation looked weak and hard to explain.",
            "Pipeline was still being tuned, so claims were not stable.",
        ],
        Inches(0.45),
        Inches(1.32),
        Inches(2.65),
        Inches(1.72),
        size=10,
    )
    add_bullets(
        slide,
        [
            "Focused on economically motivated daily returns: corn, soy, UUP, CAD.",
            "Combined 7 wheat features + 4 cross-asset returns.",
            "Used train-fold-only scaling and PCA(5), then TCN.",
        ],
        Inches(3.70),
        Inches(1.32),
        Inches(2.70),
        Inches(1.72),
        size=10,
    )
    add_bullets(
        slide,
        [
            "Small improvement over the TCN reference, but not a strong signal.",
            "Larger cross-asset + macro PCA panel underperformed on this tail.",
            "The update is better framed as an ablation, not a victory lap.",
        ],
        Inches(6.95),
        Inches(1.32),
        Inches(2.70),
        Inches(1.72),
        size=10,
    )
    add_callout(
        slide,
        "Core change: move from 'many possible alternative data sources' to a leakage-safe comparison of two PCA-augmented TCN tracks.",
        Inches(0.65),
        Inches(3.58),
        Inches(8.7),
        Inches(0.65),
        size=12,
    )
    add_textbox(
        slide,
        "Wheat features + selected cross-asset returns -> fold-safe PCA -> TCN probabilities -> holdout evaluation",
        Inches(0.90),
        Inches(4.47),
        Inches(8.2),
        Inches(0.32),
        size=13,
        bold=True,
        color=BLUE,
    )


def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "2. Table III: PCA-Augmented TCN on Terminal Holdout", 2)
    rows = [
        ["Pipeline", "Acc.", "ROC-AUC", "F1_M"],
        ["TCN_Revised reference", "0.541", "0.526", "0.541"],
        ["TCN_PCA_WHEAT_ONLY: 7 wheat + corn/soy/UUP/CAD returns, PCA(5)", "0.543", "0.506", "0.508"],
        ["TCN_PCA: Yahoo overlay + FRED-MD + 4 returns -> macro PCA(5)", "0.482", "0.501", "0.463"],
    ]
    table = add_table(slide, rows, Inches(0.42), Inches(1.03), Inches(9.15), Inches(1.52), font_size=8)
    table.table.columns[0].width = Inches(5.65)
    table.table.columns[1].width = Inches(1.05)
    table.table.columns[2].width = Inches(1.10)
    table.table.columns[3].width = Inches(1.05)
    add_callout(
        slide,
        "Key read: the narrower wheat-centered PCA track reached 54.3% accuracy, barely above the reference TCN. The broader cross-asset + macro PCA track fell below 50%.",
        Inches(0.55),
        Inches(3.00),
        Inches(4.25),
        Inches(1.05),
        size=11,
    )
    add_callout(
        slide,
        "Why it matters: alternative data helped only when it was selective and leakage-safe. More data did not automatically mean a better signal.",
        Inches(5.15),
        Inches(3.00),
        Inches(4.25),
        Inches(1.05),
        fill=LIGHT_GRAY,
        line=GRAY,
        size=11,
    )
    add_textbox(
        slide,
        "Replay: Adam LR=1e-6, 50 epochs, channels {128,64}, Apple MPS; terminal 15% chronological holdout; threshold 0.5 for Acc.",
        Inches(0.65),
        Inches(4.45),
        Inches(8.75),
        Inches(0.38),
        size=9,
        color=GRAY,
    )


def slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "3. Why MCMC? Strategy Needs Probabilities and Uncertainty", 3)
    add_bullets(
        slide,
        [
            "Investment overlay consumes P(up), not just a hard up/down label.",
            "A TCN can rank days, but a single probability does not show uncertainty.",
            "Bayesian logistic MCMC gives posterior-mean probabilities and dispersion.",
            "This supplements the strategy layer; it is not the champion forecaster.",
        ],
        Inches(0.45),
        Inches(0.92),
        Inches(4.25),
        Inches(1.95),
        size=11,
    )
    rows = [
        ["Model", "Acc.", "AUC", "Brier"],
        ["MLE logit", "0.500", "0.509", "0.255"],
        ["Bayes mean p", "0.542", "0.508", "0.249"],
    ]
    add_table(slide, rows, Inches(0.62), Inches(3.12), Inches(3.70), Inches(0.92), font_size=8)
    add_callout(
        slide,
        "MCMC role: expose calibration and uncertainty for threshold-aware long-flat rules.",
        Inches(0.62),
        Inches(4.32),
        Inches(3.85),
        Inches(0.60),
        size=10,
    )
    roc = FIG_DIR / "bayes_mcmc_roc_holdout.png"
    cal = FIG_DIR / "bayes_mcmc_calibration_holdout.png"
    if roc.exists():
        slide.shapes.add_picture(str(roc), Inches(5.05), Inches(0.92), width=Inches(2.10))
    if cal.exists():
        slide.shapes.add_picture(str(cal), Inches(7.30), Inches(0.92), width=Inches(2.10))
    unc = FIG_DIR / "bayes_mcmc_predictive_uncertainty_holdout.png"
    if unc.exists():
        slide.shapes.add_picture(str(unc), Inches(5.30), Inches(3.20), width=Inches(3.55))
    add_textbox(
        slide,
        "Default design: same wheat + cross-asset PCA inputs as TCN_PCA_WHEAT_ONLY, flattened into a linear Bayesian logit.",
        Inches(5.05),
        Inches(4.82),
        Inches(4.45),
        Inches(0.32),
        size=8,
        color=GRAY,
    )


def slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "4. Investment Strategy: Tactical Long-Flat Holdout Proxy", 4)
    pipeline = FIG_DIR / "pipeline_tcn_pca_mcmc_strategy.png"
    if pipeline.exists():
        slide.shapes.add_picture(str(pipeline), Inches(0.25), Inches(0.86), width=Inches(3.05))
    add_bullets(
        slide,
        [
            "Retrain TCN_PCA_WHEAT_ONLY on train+validation only.",
            "Score the terminal 15% holdout with predicted P(up).",
            "Go long only when P(up) clears a hurdle; otherwise stay flat.",
            "Proxy only: no commissions, rolls, margin, options, slippage, or live order timing.",
        ],
        Inches(3.50),
        Inches(0.90),
        Inches(5.95),
        Inches(1.45),
        size=10,
    )
    rows = [
        ["Rule", "Tot. ret.", "Vol.", "Sharpe", "Max DD", "% long"],
        ["Buy-and-hold", "-0.468", "0.319", "-0.36", "-0.613", "1.00"],
        ["P >= 0.52", "0.000", "0.000", "--", "0.000", "0.00"],
        ["P >= 0.50", "0.009", "0.005", "0.51", "0.000", "0.001"],
        ["P >= tau", "0.000", "0.217", "0.11", "-0.400", "0.42"],
        ["Momentum", "-0.155", "0.211", "-0.11", "-0.399", "0.46"],
    ]
    add_table(slide, rows, Inches(3.50), Inches(2.72), Inches(5.95), Inches(1.52), font_size=7)
    add_callout(
        slide,
        "Interpretation: on one bearish holdout tail, sitting out part of the decline reduced drawdown versus full exposure. This is risk management evidence, not proof of alpha.",
        Inches(3.55),
        Inches(4.48),
        Inches(5.85),
        Inches(0.62),
        size=9,
    )


def slide5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "5. Conclusion and Future Work", 5)
    add_label(slide, "What I learned", Inches(0.55), Inches(0.92), Inches(2.15), Inches(0.35), BLUE)
    add_bullets(
        slide,
        [
            "Daily wheat direction is a weak-signal problem.",
            "Chronological discipline matters more than model complexity.",
            "Alternative data must be selective, causal, and tested as ablations.",
            "AI/ML helped organize experiments, but did not remove market noise.",
        ],
        Inches(0.62),
        Inches(1.40),
        Inches(4.10),
        Inches(1.62),
        size=10,
    )
    add_label(slide, "Final conclusion", Inches(5.05), Inches(0.92), Inches(2.15), Inches(0.35), GREEN)
    add_bullets(
        slide,
        [
            "TCN and PCA-augmented TCN showed the best descriptive holdout results.",
            "Model differences should not be read as statistically significant superiority.",
            "Probability use, calibration, and when to stay flat may matter as much as accuracy.",
        ],
        Inches(5.10),
        Inches(1.40),
        Inches(4.20),
        Inches(1.62),
        size=10,
    )
    add_label(slide, "If more time were allowed", Inches(0.55), Inches(3.35), Inches(2.95), Inches(0.35), BLUE)
    add_bullets(
        slide,
        [
            "Run expanding-window walk-forward tests with nested tuning.",
            "Add FRED-vintage sensitivity and alternate calendar cutoffs.",
            "Model commissions, rolls, contract multipliers, margin, and liquidity.",
            "Compare TCN with/without macro PCA on identical splits.",
            "Use Bayesian uncertainty for calibration-aware thresholds, not overfitting.",
        ],
        Inches(0.62),
        Inches(3.82),
        Inches(8.65),
        Inches(1.05),
        size=10,
    )


def write_notes() -> None:
    NOTES_OUT.write_text(
        """# Speaker Notes: VIP Final Presentation

## Slide 1: Alternative Data Update

Last time, the alternative-data presentation did not land well because the story was too broad and still looked like tuning. The update is that I narrowed the alternative data to economically motivated daily returns: corn, soybeans, UUP as a dollar proxy, and CAD as a North American export/FX channel. I then treated this as a leakage-safe ablation: seven wheat-derived features plus four cross-asset returns, standardized and reduced by PCA only inside the training fold, then fed into the TCN.

## Slide 2: Table III

This slide is the updated alternative-data result. The reference TCN had 54.1 percent holdout accuracy. The narrower wheat-centered PCA version reached 54.3 percent, which is a very small descriptive improvement. The broader cross-asset plus macro PCA version underperformed at 48.2 percent. My interpretation is that more data is not automatically better. The useful progress is a cleaner test design and a more cautious conclusion.

## Slide 3: Why MCMC

The motivation for MCMC came from the investment strategy question. A trading overlay does not only need a label; it needs probabilities and some sense of uncertainty. The Bayesian logistic model uses the same wheat-PCA design as the PCA TCN path, but flattens the sequence into a linear probabilistic model. I do not present it as the best forecasting model. I use it to supplement the strategy layer with posterior-mean probabilities, calibration views, and uncertainty diagnostics.

## Slide 4: Investment Strategy

The strategy is intentionally simple. I retrain the TCN_PCA_WHEAT_ONLY model on train plus validation only, score the terminal holdout, and turn the predicted probability into long-flat weights. Buy-and-hold was always long, while the model rules could stay in cash. On this bearish holdout tail, buy-and-hold lost about 46.8 percent. Some long-flat rules reduced drawdown by avoiding exposure, especially the validation-threshold rule, which was long about 42 percent of the time. This is a frictionless proxy, not a full futures backtest.

## Slide 5: Conclusion and Future Work

The main lesson is forecast realism. Daily wheat direction is hard to predict, and honest chronological evaluation keeps the claims grounded. Alternative data helped only modestly and only in the narrower version. If more time were allowed, I would prioritize walk-forward tests, FRED-vintage sensitivity, broker-realistic costs and rolls, and clearer feature ablations. The final takeaway is that the project became less about claiming alpha and more about building a disciplined, reproducible ML forecasting evaluation.
""",
        encoding="utf-8",
    )


def capitalize_words_for_ppt(text: str) -> str:
    """Capitalize word starts while preserving acronyms/code-like tokens."""
    out = []
    capitalize_next = True
    for ch in text:
        if ch.isalpha() and capitalize_next:
            out.append(ch.upper())
            capitalize_next = False
        else:
            out.append(ch)
            if ch.isalpha():
                capitalize_next = False
        if ch in " \t\n\r-/:;,.()[]{}":
            capitalize_next = True
    return "".join(out)


def capitalize_slide_text(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = capitalize_words_for_ppt(run.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.text = capitalize_words_for_ppt(run.text)


def main() -> None:
    prs = Presentation(TEMPLATE)
    clear_template_slides(prs)
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    capitalize_slide_text(prs)
    prs.save(OUT)
    write_notes()
    print(f"Wrote {OUT}")
    print(f"Wrote {NOTES_OUT}")


if __name__ == "__main__":
    main()
